from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import matplotlib.pyplot as plt
import matplotlib.patches as patches

# Cria uma nova apresentação
prs = Presentation()


def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Layout com título e conteúdo
    slide = prs.slides.add_slide(slide_layout)

    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content


def add_image_slide(title, image_path):
    slide_layout = prs.slide_layouts[6]  # Layout em branco
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    left = Inches(1)
    top = Inches(1)
    pic = slide.shapes.add_picture(image_path, left, top, height=Inches(5.5))


def create_fluxograma():
    fig, ax = plt.subplots(figsize=(10, 6))

    # Adiciona retângulos para o fluxograma
    rectangles = [
        (0.1, 0.8, 'Inicializar População'),
        (0.1, 0.6, 'Avaliar Fitness da População'),
        (0.1, 0.4, 'Critério de Parada?'),
        (0.1, 0.3, 'Selecionar Indivíduos (Seleção)'),
        (0.1, 0.2, 'Cruzamento (Crossover)'),
        (0.1, 0.1, 'Mutação'),
        (0.1, 0.0, 'Nova População Gerada')
    ]

    for x, y, text in rectangles:
        rect = patches.FancyBboxPatch((x, y), 0.4, 0.1, boxstyle="round,pad=0.05", edgecolor='black',
                                      facecolor='lightblue', lw=2)
        ax.add_patch(rect)
        plt.text(x + 0.2, y + 0.05, text, ha='center', va='center', fontsize=12, fontweight='bold')

    # Adiciona setas
    arrows = [
        ((0.3, 0.8), (0.3, 0.6)),
        ((0.3, 0.6), (0.3, 0.4)),
        ((0.3, 0.4), (0.3, 0.3)),
        ((0.3, 0.3), (0.3, 0.2)),
        ((0.3, 0.2), (0.3, 0.1)),
        ((0.3, 0.1), (0.3, 0.0)),
        ((0.3, 0.4), (0.2, 0.3)),
        ((0.2, 0.3), (0.4, 0.3)),
        ((0.4, 0.3), (0.3, 0.4))
    ]

    for start, end in arrows:
        ax.annotate("", xy=end, xytext=start, arrowprops=dict(arrowstyle="->", color='black'))

    plt.xlim(0, 0.5)
    plt.ylim(-0.1, 0.9)
    plt.axis('off')
    plt.savefig('fluxograma.png', bbox_inches='tight')
    plt.close()


# Adiciona slides à apresentação
add_slide("Algoritmos Genéticos - Operadores e Exemplo",
          "Neste exemplo, vamos demonstrar o funcionamento básico de um algoritmo genético usando uma função simples de maximização \( f(x) = x^2 \). A população é composta de indivíduos representados por cadeias binárias, e aplicaremos os operadores de **seleção**, **cruzamento**, **mutação** e **elitismo**.")

add_slide("Função de Avaliação (Fitness Function)",
          "Para cada indivíduo da população, calculamos o valor da função de aptidão (fitness) como o quadrado de seu valor decimal.\n\nExemplo de uma população inicial com 4 indivíduos:\n\nIndivíduo | Binário | Decimal (x) | Fitness (f(x))\n--------- | ------- | ----------- | ---------------\nI1       | 10101   | 21          | 441\nI2       | 01100   | 12          | 144\nI3       | 11011   | 27          | 729\nI4       | 00101   | 5           | 25")

add_slide("Seleção (Roulette Wheel Selection)",
          "A seleção é realizada com base na probabilidade proporcional ao valor de fitness.\n\nNeste exemplo, escolhemos os indivíduos \( I_1 \) e \( I_3 \) para reprodução.\n\nA probabilidade de seleção é calculada como:\n\nP(I_1) = 441 / 1339 ≈ 0.33\nP(I_2) = 144 / 1339 ≈ 0.11\nP(I_3) = 729 / 1339 ≈ 0.54\nP(I_4) = 25 / 1339 ≈ 0.02")

add_slide("Cruzamento (Crossover)",
          "Realizamos o cruzamento entre \( I_1 \) e \( I_3 \) em um ponto, gerando dois novos indivíduos:\n\n- Filho 1: \(10111\) (Decimal: 23)\n- Filho 2: \(11001\) (Decimal: 25)")

add_slide("Mutação (Mutation)",
          "Aplicamos uma mutação ao último bit de \( Filho 1 \), alterando o bit de 1 para 0:\n\n10111 → 10110\n\nAgora, \( Filho 1 \) tem o valor decimal 22.")

add_slide("Elitismo",
          "Preservamos o melhor indivíduo da geração anterior \( I_3 \) com \( x = 27 \) e fitness 729.\n\nA nova população é:\n\nIndivíduo | Binário | Decimal (x) | Fitness (f(x))\n--------- | ------- | ----------- | ---------------\nI3       | 11011   | 27          | 729\nFilho 1  | 10110   | 22          | 484\nFilho 2  | 11001   | 25          | 625\nI1       | 10101   | 21          | 441")

create_fluxograma()
add_image_slide("Fluxograma - Algoritmo Genético", "../censo/fluxograma.png")

# Salva a apresentação
prs.save('algoritmos_geneticos.pptx')
